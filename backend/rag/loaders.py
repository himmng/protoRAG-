"""Per-extension document loaders and the shared text splitter."""

import os

from langchain_text_splitters import RecursiveCharacterTextSplitter


def load_document(file_path: str, filename: str):
    from langchain_core.documents import Document

    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        try:
            from langchain_community.document_loaders import PyPDFLoader
            return PyPDFLoader(file_path).load()
        except ImportError:
            raise RuntimeError("pypdf required. Run: pip install pypdf")

    if ext in (".docx", ".doc"):
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            return Docx2txtLoader(file_path).load()
        except ImportError:
            raise RuntimeError("docx2txt required. Run: pip install docx2txt")

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    slides.append(f"Slide {i}:\n" + "\n".join(texts))
            content = "\n\n".join(slides) or "(empty presentation)"
            return [Document(page_content=content, metadata={"source": file_path})]
        except ImportError:
            raise RuntimeError("python-pptx required. Run: pip install python-pptx")

    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = [
                    "\t".join("" if c is None else str(c) for c in row)
                    for row in ws.iter_rows(values_only=True)
                    if any(c is not None for c in row)
                ]
                if rows:
                    sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
            content = "\n\n".join(sheets) or "(empty workbook)"
            return [Document(page_content=content, metadata={"source": file_path})]
        except ImportError:
            raise RuntimeError("openpyxl required. Run: pip install openpyxl")

    # All remaining text-based formats: txt, md, csv, json, jsonl, yaml, xml, html, etc.
    from langchain_community.document_loaders import TextLoader
    return TextLoader(file_path, autodetect_encoding=True).load()


def get_text_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
