import asyncio
import os
import json
import re
import shutil
import traceback
from collections import defaultdict
from typing import Optional

import chromadb
from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Query
from fastapi.responses import StreamingResponse, HTMLResponse, FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

from langchain_chroma import Chroma
from langchain_openai import OpenAIEmbeddings, ChatOpenAI
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.messages import HumanMessage, SystemMessage, AIMessage

app = FastAPI(title="protoRAG+ API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

DATA_DIR = os.environ.get("DEFAULT_DATA_DIR", "./data")
MAX_HISTORY_ENTRIES = 200
MAX_UPLOAD_BYTES = 50 * 1024 * 1024
ALLOWED_EXTENSIONS = {
    ".pdf",
    ".txt", ".md", ".rst", ".log",
    ".csv", ".tsv",
    ".json", ".jsonl",
    ".yaml", ".yml",
    ".xml", ".html", ".htm",
    ".docx", ".doc",
    ".pptx", ".ppt",
    ".xlsx", ".xls",
}
_UUID_RE = re.compile(r"^[0-9a-f]{8}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{4}-[0-9a-f]{12}$")
_OFFICE_EXTS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls"}
_PREVIEW_DIRNAME = ".previews"
_session_locks: dict[str, asyncio.Lock] = defaultdict(asyncio.Lock)
# session_id → (embedding_model_name, db_path, Chroma). Reusing one Chroma per
# session across requests keeps the underlying ChromaDB system alive and the
# in-memory HNSW segment stable. Without this, every request constructs a new
# wrapper on top of the same shared system; the resulting churn races with HNSW
# writes (sync_threshold=1000 by default — small collections live in RAM),
# leading to 'Error finding id' on the next query.
# db_path is part of the cache key so a mid-session change to the data
# directory invalidates the cached wrapper instead of writing to the old path.
_session_vectorstores: dict[str, tuple[str, str, "Chroma"]] = {}


# ── Validation & path helpers ─────────────────────────────────────────────────

def _validate_session(session_id: str):
    if not _UUID_RE.match(session_id):
        raise HTTPException(status_code=400, detail="Invalid session ID")


def safe_join(base: str, *paths: str) -> str:
    base_real = os.path.realpath(base)
    candidate = os.path.realpath(os.path.join(base, *paths))
    if not candidate.startswith(base_real + os.sep) and candidate != base_real:
        raise HTTPException(status_code=400, detail="Invalid path")
    return candidate


def _effective_data_dir(custom: Optional[str]) -> str:
    """Return the data dir to use: custom override (expanded ~) or DATA_DIR default."""
    if custom and custom.strip():
        return os.path.expanduser(custom.strip())
    return DATA_DIR


def resolve_dirs(data_dir: Optional[str] = None):
    base = _effective_data_dir(data_dir)
    db_dir = os.path.join(base, "db", "session")
    docs_dir = os.path.join(base, "documents", "session")
    try:
        os.makedirs(db_dir, exist_ok=True)
        os.makedirs(docs_dir, exist_ok=True)
    except OSError as e:
        raise HTTPException(
            status_code=400,
            detail=f"Cannot create data directory '{base}': {e}",
        )
    return db_dir, docs_dir


def collection_name_for(session_id: str) -> str:
    return f"col_{session_id}"


# ── Provider / config ─────────────────────────────────────────────────────────

class ChatConfig(BaseModel):
    provider: str
    base_url: str
    api_key: str
    model_name: str
    embedding_model: str


class ChatRequest(BaseModel):
    session_id: str
    message: str
    config: ChatConfig
    data_dir: Optional[str] = None


def _safe_api_key(key: str) -> str:
    """Return 'dummy' for blank/none keys (OpenAI-compat servers ignore it anyway)."""
    return key.strip() if key and key.strip().lower() not in ("", "none") else "dummy"


def normalise_base_url(provider: str, base_url: str) -> str:
    url = base_url.rstrip("/")
    if url.endswith("/v1"):
        return url
    if provider.lower() in ("ollama", "lmstudio"):
        return url + "/v1"
    return url


def get_embeddings(config: ChatConfig) -> OpenAIEmbeddings:
    """
    All providers use OpenAI-compatible embeddings.
    For Anthropic (which has no embedding API), base_url must point to a local
    embedding service, e.g. http://localhost:11434/v1 (Ollama).
    """
    if config.provider.lower() == "anthropic":
        embed_url = config.base_url.rstrip("/")
        embed_key = "dummy"
    else:
        embed_url = normalise_base_url(config.provider, config.base_url)
        embed_key = _safe_api_key(config.api_key)

    return OpenAIEmbeddings(
        model=config.embedding_model,
        openai_api_base=embed_url,
        openai_api_key=embed_key,
        check_embedding_ctx_length=False,
    )


def get_llm(config: ChatConfig):
    provider = config.provider.lower()

    if provider == "anthropic":
        try:
            from langchain_anthropic import ChatAnthropic
        except ImportError:
            raise RuntimeError(
                "langchain-anthropic not installed. Run: pip install langchain-anthropic"
            )
        return ChatAnthropic(
            model=config.model_name,
            anthropic_api_key=config.api_key,
            streaming=True,
            max_tokens=8192,
        )

    # All OpenAI-compatible providers: Ollama, LM Studio, LiteLLM, OpenAI, Custom
    base_url = normalise_base_url(config.provider, config.base_url)
    api_key = _safe_api_key(config.api_key)
    return ChatOpenAI(
        model=config.model_name,
        openai_api_base=base_url,
        openai_api_key=api_key,
        streaming=True,
    )


# ── Background index rebuild ──────────────────────────────────────────────────

async def _rebuild_index_background(
    session_id: str, db_path: str, docs_dir: str, db_dir: str, config: "ChatConfig"
):
    """Re-index all session documents after a corruption event.

    The corrupted system is already released by chat()'s exception handler via
    _release_chroma_system. We call it again here as belt-and-suspenders before
    the wipe — guarantees no cached system holds open handles on the directory
    we're about to delete.
    """
    session_doc_path = os.path.join(docs_dir, session_id)
    doc_files = sorted(
        f for f in (os.listdir(session_doc_path) if os.path.exists(session_doc_path) else [])
        if os.path.isfile(os.path.join(session_doc_path, f))
    )

    async with _session_locks[session_id]:
        _release_session(session_id, db_path)
        if os.path.exists(db_path):
            try:
                shutil.rmtree(db_path)
                print(f"[RAG] Wiped corrupted ChromaDB for session '{session_id}'")
            except Exception as e:
                print(f"[RAG] Could not wipe ChromaDB: {e}")

        meta_path = os.path.join(db_dir, session_id, "embed_meta.json")

        if not doc_files:
            try:
                if os.path.exists(meta_path):
                    with open(meta_path) as f:
                        _m = json.load(f)
                    _m.pop("corrupted", None)
                    with open(meta_path, "w") as f:
                        json.dump(_m, f)
            except Exception:
                pass
            return

        rebuilt: list[str] = []
        try:
            embeddings = get_embeddings(config)
            vs = get_vectorstore(db_path, embeddings, session_id)
            for doc_file in doc_files:
                try:
                    rdocs = load_document(os.path.join(session_doc_path, doc_file), doc_file)
                    rsplits = get_text_splitter().split_documents(rdocs)
                    for chunk in rsplits:
                        chunk.metadata["doc_filename"] = doc_file
                    vs.add_documents(rsplits)
                    rebuilt.append(doc_file)
                    print(f"[RAG] Background-rebuilt '{doc_file}' ({len(rsplits)} chunks)")
                except Exception as e:
                    print(f"[RAG] Background-rebuild failed for '{doc_file}': {e}")
        except Exception as e:
            print(f"[RAG] Background-rebuild aborted: {e}")
            try:
                if os.path.exists(meta_path):
                    with open(meta_path) as f:
                        _m = json.load(f)
                    _m.pop("corrupted", None)
                    with open(meta_path, "w") as f:
                        json.dump(_m, f)
            except Exception:
                pass
            return

        try:
            existing_meta: dict = {}
            if os.path.exists(meta_path):
                with open(meta_path) as f:
                    existing_meta = json.load(f)
            existing_meta["indexed_files"] = sorted(rebuilt)
            existing_meta.pop("corrupted", None)
            os.makedirs(os.path.dirname(meta_path), exist_ok=True)
            with open(meta_path, "w") as f:
                json.dump(existing_meta, f)
        except Exception:
            pass
        print(f"[RAG] Background rebuild complete for session '{session_id}': {rebuilt}")


# ── ChromaDB helpers ──────────────────────────────────────────────────────────

def _chroma_client(db_path: str) -> chromadb.PersistentClient:
    os.makedirs(db_path, exist_ok=True)
    return chromadb.PersistentClient(path=db_path)


def _release_chroma_system(db_path: str) -> None:
    """Force-release the cached ChromaDB system for this path.

    ChromaDB's SharedSystemClient caches systems by the exact persist_directory
    string passed to PersistentClient. Re-opening the same path returns the
    cached system — including its open SQLite/HNSW handles. When the underlying
    HNSW index becomes corrupted ('Error finding id'), we must explicitly stop
    that cached system before wiping the directory; otherwise the next
    PersistentClient hands back the dead system, whose stale handles cause
    SQLITE_READONLY_DIRECTORY (code 1032) on the first write to the rebuilt db.
    """
    try:
        from chromadb.api.shared_system_client import SharedSystemClient
        with SharedSystemClient._refcount_lock:
            system = SharedSystemClient._identifier_to_system.pop(db_path, None)
            SharedSystemClient._identifier_to_refcount.pop(db_path, None)
        if system is not None:
            system.stop()
            print(f"[RAG] Released cached ChromaDB system for '{db_path}'")
    except Exception as e:
        print(f"[RAG] Could not release ChromaDB system for '{db_path}': {e}")


def _release_session(session_id: str, db_path: str) -> None:
    """Drop the cached Chroma for this session and force-release the system.

    Call before any operation that wipes or recreates the on-disk database
    (corruption recovery, re-upload, session/document delete) — otherwise the
    next get_vectorstore() returns a Chroma whose handles point at the now-
    deleted directory.
    """
    cached = _session_vectorstores.pop(session_id, None)
    if cached is not None:
        try:
            cached[2]._client.close()
        except Exception as e:
            print(f"[RAG] Error closing cached Chroma for '{session_id}': {e}")
        # Also release the system at the cached path in case it differs
        # from db_path (data dir was changed mid-session).
        if cached[1] != db_path:
            _release_chroma_system(cached[1])
    _release_chroma_system(db_path)


def get_vectorstore(db_path: str, embeddings, session_id: str) -> Chroma:
    """Return a session-cached Chroma, reusing one wrapper across requests.

    The cache is invalidated when either the embedding model or the on-disk
    db_path changes — the embedding function and persist_directory are bound
    at Chroma construction, so a stale wrapper would silently embed new uploads
    with the wrong model or write to the wrong directory.
    """
    model = getattr(embeddings, "model", "") or ""
    cached = _session_vectorstores.get(session_id)
    if cached is not None and cached[0] == model and cached[1] == db_path:
        return cached[2]

    if cached is not None:
        try:
            cached[2]._client.close()
        except Exception:
            pass
        if cached[1] != db_path:
            _release_chroma_system(cached[1])
        _session_vectorstores.pop(session_id, None)

    os.makedirs(db_path, exist_ok=True)
    vs = Chroma(
        persist_directory=db_path,
        embedding_function=embeddings,
        collection_name=collection_name_for(session_id),
    )
    _session_vectorstores[session_id] = (model, db_path, vs)
    return vs


def _delete_vectors_for_file(db_path: str, session_id: str, filename: str) -> int:
    """
    Delete all vectors tagged with doc_filename==filename using the raw ChromaDB
    client — no embedding function required, so no dimension-mismatch errors.
    """
    if not os.path.exists(db_path):
        return 0
    try:
        client = _chroma_client(db_path)
        col_name = collection_name_for(session_id)
        try:
            collection = client.get_collection(col_name)
        except Exception:
            return 0  # collection doesn't exist yet

        result = collection.get(where={"doc_filename": filename})
        ids = result.get("ids", [])
        if ids:
            collection.delete(ids=ids)
            print(f"[RAG] Deleted {len(ids)} vectors for '{filename}'")
        return len(ids)
    except Exception as e:
        print(f"[RAG] Warning: could not delete vectors for '{filename}': {e}")
        traceback.print_exc()
        return 0


def _count_vectors(db_path: str, session_id: str) -> int:
    if not os.path.exists(db_path):
        return 0
    try:
        collection = _chroma_client(db_path).get_collection(collection_name_for(session_id))
        return collection.count()
    except Exception:
        return 0


# ── Document loading ──────────────────────────────────────────────────────────

def load_document(file_path: str, filename: str):
    from langchain_core.documents import Document

    ext = os.path.splitext(filename.lower())[1]

    if ext == ".pdf":
        try:
            from langchain_community.document_loaders import PyPDFLoader
            return PyPDFLoader(file_path).load()
        except ImportError:
            raise RuntimeError("pypdf required. Run: pip install pypdf")

    if ext in (".docx", ".doc"):
        try:
            from langchain_community.document_loaders import Docx2txtLoader
            return Docx2txtLoader(file_path).load()
        except ImportError:
            raise RuntimeError("docx2txt required. Run: pip install docx2txt")

    if ext in (".pptx", ".ppt"):
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [
                    shape.text.strip()
                    for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                if texts:
                    slides.append(f"Slide {i}:\n" + "\n".join(texts))
            content = "\n\n".join(slides) or "(empty presentation)"
            return [Document(page_content=content, metadata={"source": file_path})]
        except ImportError:
            raise RuntimeError("python-pptx required. Run: pip install python-pptx")

    if ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path, data_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = [
                    "\t".join("" if c is None else str(c) for c in row)
                    for row in ws.iter_rows(values_only=True)
                    if any(c is not None for c in row)
                ]
                if rows:
                    sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
            content = "\n\n".join(sheets) or "(empty workbook)"
            return [Document(page_content=content, metadata={"source": file_path})]
        except ImportError:
            raise RuntimeError("openpyxl required. Run: pip install openpyxl")

    # All remaining text-based formats: txt, md, csv, json, jsonl, yaml, xml, html, etc.
    from langchain_community.document_loaders import TextLoader
    return TextLoader(file_path, autodetect_encoding=True).load()


def get_text_splitter() -> RecursiveCharacterTextSplitter:
    return RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)


# ── History helpers ───────────────────────────────────────────────────────────

def load_history(session_id: str, db_dir: str) -> list:
    path = os.path.join(db_dir, session_id, "history.json")
    if os.path.exists(path):
        with open(path) as f:
            return json.load(f)
    return []


def save_history(session_id: str, history: list, db_dir: str):
    path = os.path.join(db_dir, session_id, "history.json")
    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "w") as f:
        json.dump(history[-MAX_HISTORY_ENTRIES:], f)


# ── Upload ────────────────────────────────────────────────────────────────────

@app.post("/api/upload")
async def upload_document(
    session_id: str = Form(...),
    provider: str = Form(...),
    base_url: str = Form(...),
    api_key: str = Form(...),
    model_name: str = Form(...),
    embedding_model: str = Form(...),
    data_dir: Optional[str] = Form(None),
    file: UploadFile = File(...),
):
    try:
        _validate_session(session_id)

        ext = os.path.splitext(file.filename or "")[1].lower()
        if ext not in ALLOWED_EXTENSIONS:
            return {"status": "error", "message": f"Unsupported file type '{ext}'. Allowed: {', '.join(sorted(ALLOWED_EXTENSIONS))}"}

        content = await file.read()
        if len(content) > MAX_UPLOAD_BYTES:
            return {"status": "error", "message": f"File too large (max {MAX_UPLOAD_BYTES // 1024 // 1024} MB)."}

        config = ChatConfig(
            provider=provider, base_url=base_url, api_key=api_key,
            model_name=model_name, embedding_model=embedding_model,
        )

        db_dir, docs_dir = resolve_dirs(data_dir)
        session_doc_path = safe_join(docs_dir, session_id)
        os.makedirs(session_doc_path, exist_ok=True)

        filename = os.path.basename(file.filename or "upload")
        file_path = safe_join(session_doc_path, filename)
        with open(file_path, "wb") as f:
            f.write(content)

        docs = load_document(file_path, filename)
        splits = get_text_splitter().split_documents(docs)
        for chunk in splits:
            chunk.metadata["doc_filename"] = filename

        if not splits:
            return {"status": "warning", "message": f"'{filename}' was saved but produced 0 text chunks — it may be empty, image-only, or unsupported internally. RAG will not work for this file."}

        db_path = safe_join(db_dir, session_id, "chromadb")

        async with _session_locks[session_id]:
            embeddings = get_embeddings(config)

            # Read existing meta (tracks which filenames are already indexed).
            # Using a plain JSON file avoids opening a second ChromaDB client to
            # do the re-upload check — dual PersistentClient instances on the same
            # path cause SQLite WAL conflicts that crash the upsert.
            meta_path = os.path.join(db_dir, session_id, "embed_meta.json")
            meta: dict = {}
            if os.path.exists(meta_path):
                try:
                    with open(meta_path) as f:
                        meta = json.load(f)
                except Exception:
                    pass

            indexed_files: set = set(meta.get("indexed_files", []))
            is_reupload = filename in indexed_files

            if is_reupload:
                # Wipe the ChromaDB directory entirely so the rebuild starts from a
                # clean slate — no ghost HNSW entries, no lock conflicts. Release
                # the cached Chroma first so its handles aren't pointing at a
                # deleted dir when get_vectorstore() reopens below.
                _release_session(session_id, db_path)
                if os.path.exists(db_path):
                    shutil.rmtree(db_path)
                    print(f"[RAG] Wiped ChromaDB for rebuild (re-upload of '{filename}')")

            # Single LangChain client for all indexing in this request.
            vectorstore = get_vectorstore(db_path, embeddings, session_id)
            vectorstore.add_documents(splits)
            indexed_files.add(filename)

            if is_reupload:
                for other in os.listdir(session_doc_path):
                    if other == filename or not os.path.isfile(os.path.join(session_doc_path, other)):
                        continue
                    try:
                        other_docs = load_document(os.path.join(session_doc_path, other), other)
                        other_splits = get_text_splitter().split_documents(other_docs)
                        for chunk in other_splits:
                            chunk.metadata["doc_filename"] = other
                        vectorstore.add_documents(other_splits)
                        indexed_files.add(other)
                        print(f"[RAG] Re-indexed '{other}' during rebuild")
                    except Exception as e:
                        print(f"[RAG] Could not re-index '{other}': {e}")

            print(f"[RAG] Indexed {len(splits)} chunks for '{filename}'")

            os.makedirs(os.path.dirname(meta_path), exist_ok=True)
            with open(meta_path, "w") as f:
                json.dump({
                    "embedding_model": config.embedding_model,
                    "provider": config.provider,
                    "indexed_files": sorted(indexed_files),
                }, f)

        return {"status": "success", "message": f"Indexed {filename} ({len(splits)} chunks)."}

    except Exception as e:
        traceback.print_exc()
        return {"status": "error", "message": str(e)}


# ── @mention parsing ──────────────────────────────────────────────────────────

def parse_at_mentions(message: str) -> tuple[list[str], str]:
    pattern = r'@"([^"]+)"|@(\S+)'
    mentions = [m.group(1) or m.group(2) for m in re.finditer(pattern, message)]
    cleaned = re.sub(pattern, "", message).strip()
    return mentions, cleaned


# ── Chat ──────────────────────────────────────────────────────────────────────

@app.post("/api/chat")
async def chat(request: ChatRequest):
    _validate_session(request.session_id)
    session_id = request.session_id
    raw_message = request.message
    config = request.config

    db_dir, docs_dir = resolve_dirs(request.data_dir)
    mentioned_files, query = parse_at_mentions(raw_message)
    if not query:
        query = raw_message

    db_path = os.path.join(db_dir, session_id, "chromadb")
    context = ""
    retrieval_warning = ""

    if os.path.exists(db_path):
        try:
            meta_path = os.path.join(db_dir, session_id, "embed_meta.json")
            if os.path.exists(meta_path):
                with open(meta_path) as f:
                    meta = json.load(f)
                if meta.get("embedding_model") != config.embedding_model:
                    retrieval_warning = (
                        f"\n\n> ⚠️ **RAG skipped — embedding model mismatch**: "
                        f"documents indexed with `{meta['embedding_model']}` "
                        f"but settings use `{config.embedding_model}`. "
                        "Re-upload your documents to fix this."
                    )
                    raise ValueError("embedding model mismatch")

                if meta.get("corrupted"):
                    retrieval_warning = (
                        "\n\n> ⚠️ **RAG index is being rebuilt** — "
                        "please ask your question again in a moment."
                    )
                    raise ValueError("index being rebuilt")

            embeddings = get_embeddings(config)
            vectorstore = get_vectorstore(db_path, embeddings, session_id)

            # Count through the already-open LangChain client — avoids opening a
            # second raw PersistentClient whose SQLite WAL snapshot may not include
            # vectors written by the LangChain client, breaking incremental uploads.
            total = vectorstore._collection.count()
            if total == 0:
                raise ValueError("empty collection — no vectors indexed yet")

            k = min(8, total)
            fetch_k = min(30, max(k + 1, total))

            docs = []
            if mentioned_files:
                session_doc_path = os.path.join(docs_dir, session_id)
                available = set(
                    f for f in (os.listdir(session_doc_path) if os.path.exists(session_doc_path) else [])
                    if os.path.isfile(os.path.join(session_doc_path, f))
                )
                valid = [f for f in mentioned_files if f in available]
                missing = [f for f in mentioned_files if f not in available]

                if missing:
                    retrieval_warning += (
                        "\n\n> ⚠️ **Document(s) not found in this session**: "
                        + ", ".join(f"`{f}`" for f in missing)
                    )

                if valid:
                    where_filter = (
                        {"doc_filename": {"$eq": valid[0]}}
                        if len(valid) == 1
                        else {"doc_filename": {"$in": valid}}
                    )
                    # ChromaDB Rust backend crashes with MMR + where filter — use
                    # plain similarity_search which goes through a different code path.
                    docs = vectorstore.similarity_search(query, k=k, filter=where_filter)
                    print(f"[RAG] @filter {where_filter} → {len(docs)} chunks")
            else:
                docs = vectorstore.max_marginal_relevance_search(query, k=k, fetch_k=fetch_k)
                sources = list({d.metadata.get("doc_filename", "?") for d in docs})
                print(f"[RAG] MMR → {len(docs)} chunks from {sources}")

            if docs:
                grouped: dict[str, list[str]] = {}
                for d in docs:
                    src = d.metadata.get("doc_filename", "unknown")
                    grouped.setdefault(src, []).append(d.page_content)
                context = "\n\n---\n\n".join(
                    f"### Source: {src}\n" + "\n\n".join(chunks)
                    for src, chunks in grouped.items()
                )

        except Exception as e:
            print(f"[RAG] Retrieval error: {e}")
            traceback.print_exc()
            if "Error finding id" in str(e):
                # Mark the index as corrupted so subsequent requests skip RAG
                # while the background task wipes and rebuilds. Release the
                # cached Chroma *now* so its SQLite/HNSW handles are closed
                # before the rebuild wipes the directory — otherwise the next
                # PersistentClient hits SQLITE_READONLY_DIRECTORY (1032).
                _release_session(session_id, db_path)
                meta_path = os.path.join(db_dir, session_id, "embed_meta.json")
                try:
                    _m: dict = {}
                    if os.path.exists(meta_path):
                        with open(meta_path) as f:
                            _m = json.load(f)
                    _m["indexed_files"] = []
                    _m["corrupted"] = True
                    os.makedirs(os.path.dirname(meta_path), exist_ok=True)
                    with open(meta_path, "w") as f:
                        json.dump(_m, f)
                except Exception:
                    pass

                asyncio.create_task(
                    _rebuild_index_background(session_id, db_path, docs_dir, db_dir, config)
                )
                retrieval_warning = (
                    "\n\n> ⚠️ **RAG index was corrupted.** "
                    "Your documents are being re-indexed in the background — "
                    "please ask your question again in a moment."
                )
            elif not retrieval_warning:
                retrieval_warning = f"\n\n> ⚠️ **RAG retrieval failed**: {e}"

    history = load_history(session_id, db_dir)

    session_doc_path = os.path.join(docs_dir, session_id)
    session_docs = sorted(
        f for f in (os.listdir(session_doc_path) if os.path.exists(session_doc_path) else [])
        if os.path.isfile(os.path.join(session_doc_path, f))
    )
    doc_list_note = (
        f"The user has uploaded {len(session_docs)} document(s) to this session: "
        + ", ".join(f'"{d}"' for d in session_docs)
        + "."
    ) if session_docs else "No documents have been uploaded to this session yet."

    sys_prompt = (
        "You are a helpful and intelligent AI assistant. "
        f"{doc_list_note} "
        "When answering from document context, always specify which document your answer comes from."
    )
    messages = [SystemMessage(content=sys_prompt)]
    for msg in history[-10:]:
        role = msg.get("role", "")
        if role == "user":
            messages.append(HumanMessage(content=msg.get("content", "")))
        elif role == "assistant":
            messages.append(AIMessage(content=msg.get("content", "")))

    if context:
        source_note = f"(filtered to: {', '.join(mentioned_files)})" if mentioned_files else ""
        all_docs_line = (
            f"Session documents ({len(session_docs)}): {', '.join(session_docs)}.\n"
            "The excerpts below are the most relevant chunks retrieved for this query.\n\n"
        ) if session_docs else ""
        augmented = (
            f"{all_docs_line}"
            f"Retrieved context {source_note}:\n"
            "---------------------\n"
            f"{context}\n"
            "---------------------\n"
            "Answer using the context above and cite which document(s) your answer comes from. "
            "For questions about how many documents are in the session or their names, "
            "use the full session document list stated above — not just what was retrieved.\n\n"
            f"User Query: {query}"
        )
        messages.append(HumanMessage(content=augmented))
    else:
        messages.append(HumanMessage(content=query))

    llm = get_llm(config)
    lock = _session_locks[session_id]

    async def generate():
        full_response = ""
        try:
            async for chunk in llm.astream(messages):
                token = chunk.content if hasattr(chunk, "content") else str(chunk)
                if token:
                    full_response += token
                    yield f"data: {json.dumps({'content': token})}\n\n"
        except Exception as e:
            err = f"\n\n**Connection Error:** {e}"
            full_response += err
            yield f"data: {json.dumps({'content': err})}\n\n"

        if retrieval_warning:
            full_response += retrieval_warning
            yield f"data: {json.dumps({'content': retrieval_warning})}\n\n"

        async with lock:
            current = load_history(session_id, db_dir)
            current.append({"role": "user", "content": query})
            current.append({"role": "assistant", "content": full_response})
            save_history(session_id, current, db_dir)

        yield "data: [DONE]\n\n"

    return StreamingResponse(generate(), media_type="text/event-stream")


# ── Session management ────────────────────────────────────────────────────────

@app.get("/api/sessions")
def list_sessions(data_dir: Optional[str] = Query(None)):
    db_dir, docs_dir = resolve_dirs(data_dir)
    sessions = []

    if not os.path.exists(db_dir):
        return {"sessions": []}

    for sid in os.listdir(db_dir):
        if not os.path.isdir(os.path.join(db_dir, sid)):
            continue
        hist_path = os.path.join(db_dir, sid, "history.json")
        doc_path = os.path.join(docs_dir, sid)

        has_docs = (
            os.path.exists(doc_path)
            and any(os.path.isfile(os.path.join(doc_path, f)) for f in os.listdir(doc_path))
        ) if os.path.exists(doc_path) else False

        preview = "New Chat"
        timestamp = 0
        if os.path.exists(hist_path):
            try:
                with open(hist_path) as f:
                    hist = json.load(f)
                if hist:
                    last_user = next(
                        (m for m in reversed(hist) if m.get("role") == "user"),
                        hist[-1],
                    )
                    content = last_user.get("content", "")
                    preview = (content[:35] + "...") if len(content) > 35 else (content or "New Chat")
                timestamp = os.path.getmtime(hist_path)
            except Exception:
                preview = "(history corrupted)"

        sessions.append({"id": sid, "preview": preview, "timestamp": timestamp, "is_rag": has_docs})

    sessions.sort(key=lambda x: x["timestamp"], reverse=True)
    return {"sessions": sessions}


@app.get("/api/sessions/{session_id}")
def get_session(session_id: str, data_dir: Optional[str] = Query(None)):
    _validate_session(session_id)
    db_dir, docs_dir = resolve_dirs(data_dir)
    history = load_history(session_id, db_dir)
    doc_path = os.path.join(docs_dir, session_id)
    docs = [
        f for f in (os.listdir(doc_path) if os.path.exists(doc_path) else [])
        if os.path.isfile(os.path.join(doc_path, f))
    ]
    return {"history": history, "documents": docs}


@app.delete("/api/sessions/{session_id}")
def delete_session(session_id: str, data_dir: Optional[str] = Query(None)):
    _validate_session(session_id)
    db_dir, docs_dir = resolve_dirs(data_dir)
    db_path = os.path.join(db_dir, session_id, "chromadb")
    _release_session(session_id, db_path)
    for path in [os.path.join(db_dir, session_id), os.path.join(docs_dir, session_id)]:
        if os.path.exists(path):
            shutil.rmtree(path)
    return {"status": "success"}


def _soffice_bin() -> str | None:
    return shutil.which("soffice") or shutil.which("libreoffice")


def _preview_cache_path(session_dir: str, filename: str) -> str:
    base = os.path.splitext(filename)[0]
    previews_dir = safe_join(session_dir, _PREVIEW_DIRNAME)
    return safe_join(previews_dir, base + ".pdf")


def _clear_preview_cache_for(session_dir: str, filename: str) -> None:
    try:
        cached = _preview_cache_path(session_dir, filename)
        if os.path.exists(cached):
            os.remove(cached)
    except HTTPException:
        pass


async def _convert_office_to_pdf(src_path: str, session_dir: str, filename: str) -> str:
    soffice = _soffice_bin()
    if not soffice:
        raise HTTPException(
            status_code=503,
            detail="LibreOffice not installed on backend host. Install it to preview Office files (e.g. `sudo apt install libreoffice` or `brew install --cask libreoffice`).",
        )

    cached = _preview_cache_path(session_dir, filename)
    if os.path.exists(cached) and os.path.getmtime(cached) >= os.path.getmtime(src_path):
        return cached

    out_dir = os.path.dirname(cached)
    os.makedirs(out_dir, exist_ok=True)

    # Use a per-conversion user profile so concurrent conversions don't collide
    # on LibreOffice's single-instance lock.
    profile_dir = os.path.join(out_dir, f".lo_profile_{os.getpid()}")
    os.makedirs(profile_dir, exist_ok=True)
    proc = await asyncio.create_subprocess_exec(
        soffice,
        f"-env:UserInstallation=file://{profile_dir}",
        "--headless", "--convert-to", "pdf",
        "--outdir", out_dir, src_path,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    try:
        _, stderr = await asyncio.wait_for(proc.communicate(), timeout=90)
    except asyncio.TimeoutError:
        proc.kill()
        raise HTTPException(status_code=504, detail="LibreOffice conversion timed out")

    if proc.returncode != 0 or not os.path.exists(cached):
        msg = stderr.decode(errors="ignore").strip() or "unknown error"
        raise HTTPException(status_code=500, detail=f"LibreOffice conversion failed: {msg[:300]}")
    return cached


@app.get("/api/sessions/{session_id}/documents/{filename}/preview")
async def preview_document(session_id: str, filename: str, data_dir: Optional[str] = Query(None)):
    _validate_session(session_id)
    _, docs_dir = resolve_dirs(data_dir)
    session_dir = safe_join(docs_dir, session_id)
    safe_filename = os.path.basename(filename)
    file_path = safe_join(session_dir, safe_filename)
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")

    ext = os.path.splitext(safe_filename.lower())[1]
    if ext in _OFFICE_EXTS:
        pdf_path = await _convert_office_to_pdf(file_path, session_dir, safe_filename)
        return FileResponse(pdf_path, media_type="application/pdf")
    return FileResponse(file_path)


@app.get("/api/sessions/{session_id}/documents/{filename}")
def get_document(session_id: str, filename: str, data_dir: Optional[str] = Query(None)):
    _validate_session(session_id)
    _, docs_dir = resolve_dirs(data_dir)
    session_dir = safe_join(docs_dir, session_id)
    file_path = safe_join(session_dir, os.path.basename(filename))
    if not os.path.exists(file_path):
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(file_path)


@app.delete("/api/sessions/{session_id}/documents/{filename}")
async def delete_document(session_id: str, filename: str, data_dir: Optional[str] = Query(None)):
    _validate_session(session_id)
    db_dir, docs_dir = resolve_dirs(data_dir)

    session_dir = safe_join(docs_dir, session_id)
    safe_filename = os.path.basename(filename)
    file_path = safe_join(session_dir, safe_filename)

    if os.path.exists(file_path):
        os.remove(file_path)
        print(f"[DOC] Removed '{safe_filename}' from session '{session_id}'")

    _clear_preview_cache_for(session_dir, safe_filename)

    db_path = safe_join(db_dir, session_id, "chromadb")
    remaining = []

    async with _session_locks[session_id]:
        _delete_vectors_for_file(db_path, session_id, safe_filename)

        remaining = [
            f for f in (os.listdir(session_dir) if os.path.exists(session_dir) else [])
            if os.path.isfile(os.path.join(session_dir, f))
        ]

        if not remaining and os.path.exists(db_path):
            _release_session(session_id, db_path)
            shutil.rmtree(db_path)
            print(f"[RAG] Wiped ChromaDB for session '{session_id}' (no documents left)")

        # Remove the deleted file from the indexed_files tracking in meta.
        meta_path = os.path.join(db_dir, session_id, "embed_meta.json")
        if os.path.exists(meta_path):
            try:
                with open(meta_path) as f:
                    meta = json.load(f)
                meta["indexed_files"] = sorted(
                    x for x in meta.get("indexed_files", []) if x != safe_filename
                )
                with open(meta_path, "w") as f:
                    json.dump(meta, f)
            except Exception:
                pass

    return {"status": "success", "remaining_documents": remaining}


# ── Lifecycle ─────────────────────────────────────────────────────────────────

@app.on_event("shutdown")
def _flush_session_vectorstores():
    """Close all cached Chroma instances so HNSW state is flushed to disk.

    Without this, sessions with fewer than sync_threshold (1000) vectors live
    only in memory and are lost when the process exits — which happens on every
    reload and Ctrl+C.
    """
    for session_id, (_, _, vs) in list(_session_vectorstores.items()):
        try:
            vs._client.close()
        except Exception as e:
            print(f"[RAG] Error closing Chroma for '{session_id}' on shutdown: {e}")
    _session_vectorstores.clear()


# ── UI ────────────────────────────────────────────────────────────────────────

@app.get("/")
def get_ui():
    if os.path.exists("index.html"):
        with open("index.html", "r", encoding="utf-8") as f:
            return HTMLResponse(content=f.read())
    return HTMLResponse(content="<h1>index.html not found alongside backend.py</h1>")


if __name__ == "__main__":
    import uvicorn
    print("Starting protoRAG+...")
    print("Open http://localhost:8000 in your browser.")
    uvicorn.run("backend:app", host="0.0.0.0", port=8000, reload=True)
